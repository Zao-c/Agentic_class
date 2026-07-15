import hashlib
import re
import threading
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional

from app.config import Settings, settings as default_settings
from app.schemas import Citation
from app.schemas import RetrievalStrategy
from app.storage import Store


SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx"}
QUESTION_BANK_MARKERS = {"题库", "习题", "参考答案", "答案"}


def tokenize(text: str) -> List[str]:
    normalized = re.sub(r"\s+", "", text.lower())
    latin = re.findall(r"[a-z]+(?:[_-][a-z0-9]+)*|\d+(?:\.\d+)?", normalized)
    chinese_segments = re.findall(r"[\u4e00-\u9fff]+", normalized)
    chinese_tokens: List[str] = []
    for segment in chinese_segments:
        chinese_tokens.extend(list(segment))
        chinese_tokens.extend(segment[index : index + 2] for index in range(len(segment) - 1))
    return latin + chinese_tokens


def split_text(text: str, max_chars: int = 900, overlap: int = 120) -> List[str]:
    text = re.sub(r"\r\n?", "\n", text).strip()
    if not text:
        return []
    sections = re.split(r"(?=^#{1,4}\s|^第[一二三四五六七八九十\d]+[章节]\s*|^\d+[.、]\s*)", text, flags=re.M)
    chunks: List[str] = []
    buffer = ""
    for section in sections:
        section = section.strip()
        if not section:
            continue
        if len(buffer) + len(section) + 1 <= max_chars:
            buffer = f"{buffer}\n{section}".strip()
            continue
        if buffer:
            chunks.append(buffer)
        if len(section) <= max_chars:
            buffer = section
            continue
        start = 0
        while start < len(section):
            chunks.append(section[start : start + max_chars])
            start += max(1, max_chars - overlap)
        buffer = ""
    if buffer:
        chunks.append(buffer)
    return chunks


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    if suffix == ".docx":
        from docx import Document

        document = Document(str(path))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                paragraphs.append(" | ".join(cell.text.strip() for cell in row.cells))
        return "\n".join(paragraphs)
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(path))
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append("第%d页\n%s" % (index, "\n".join(texts)))
        return "\n\n".join(slides)
    raise ValueError("不支持的文件类型: %s" % suffix)


class Retriever:
    def __init__(self, store: Store, settings: Settings = default_settings):
        self.store = store
        self.settings = settings
        from app.ranking import LocalRetrievalEngine

        self.engine = LocalRetrievalEngine(settings)
        self._chunks_cache: Optional[List[Dict[str, Any]]] = None
        self._chunks_cache_count = -1
        self._cache_lock = threading.RLock()

    def _active_chunks(self, access_scopes: Optional[List[str]]) -> List[Dict[str, Any]]:
        if access_scopes is not None:
            return self.store.active_chunks(access_scopes)
        with self._cache_lock:
            current_count = self.store.count_chunks()
            if self._chunks_cache is None or current_count != self._chunks_cache_count:
                self._chunks_cache = self.store.active_chunks()
                self._chunks_cache_count = current_count
            return self._chunks_cache

    def prepare(self, strategy: Optional[str] = None) -> None:
        selected = RetrievalStrategy(strategy) if strategy else None
        self.engine.prepare(self._active_chunks(None), selected)

    def import_text(
        self,
        title: str,
        content: str,
        document_type: str = "course_material",
        source_path: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        cleaned = content.strip()
        if not cleaned:
            raise ValueError("文档没有可提取文本")
        content_hash = hashlib.sha256(cleaned.encode("utf-8")).hexdigest()
        existing = self.store.document_by_hash(content_hash)
        if existing:
            return {
                "document_id": existing["document_id"],
                "status": "duplicate",
                "chunks": 0,
                "content_hash": content_hash,
            }
        document_id = "doc_" + uuid.uuid4().hex
        merged_metadata = {
            "course": "工业机器人",
            "chapter": None,
            "page": None,
            "version": "1",
            "effective_date": None,
            "access_scope": "public",
        }
        merged_metadata.update(metadata or {})
        chunk_rows = []
        for index, chunk in enumerate(split_text(cleaned)):
            chunk_rows.append(
                {
                    "chunk_id": "chk_" + uuid.uuid4().hex,
                    "chunk_index": index,
                    "content": chunk,
                    "tokens": tokenize(chunk),
                    "metadata": merged_metadata,
                }
            )
        if not chunk_rows:
            raise ValueError("文档切分后没有有效片段")
        self.store.add_document(
            document_id,
            title,
            source_path,
            document_type,
            merged_metadata,
            content_hash,
            chunk_rows,
        )
        self.engine.invalidate()
        with self._cache_lock:
            self._chunks_cache = None
            self._chunks_cache_count = -1
        return {
            "document_id": document_id,
            "status": "created",
            "chunks": len(chunk_rows),
            "content_hash": content_hash,
        }

    def import_path(
        self,
        path: Path,
        document_type: str = "course_material",
        metadata: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        path = Path(path).resolve()
        content = extract_text(path)
        return self.import_text(path.stem, content, document_type, str(path), metadata)

    def import_directory(self, root: Path, include_binary: bool = False) -> Dict[str, Any]:
        root = Path(root)
        extensions = SUPPORTED_EXTENSIONS if include_binary else {".txt", ".md"}
        created = duplicate = failed = 0
        errors = []
        for path in sorted(root.rglob("*")):
            if not path.is_file() or path.suffix.lower() not in extensions:
                continue
            relative = path.relative_to(root)
            document_type = relative.parts[0] if len(relative.parts) > 1 else "course_material"
            try:
                result = self.import_path(
                    path,
                    document_type=document_type,
                    metadata={"source_relative_path": str(relative).replace("\\", "/")},
                )
                if result["status"] == "created":
                    created += 1
                else:
                    duplicate += 1
            except Exception as exc:  # keep one broken document from stopping ingestion
                failed += 1
                errors.append({"path": str(relative), "error": str(exc)})
        return {"created": created, "duplicate": duplicate, "failed": failed, "errors": errors}

    def search(
        self,
        query: str,
        top_k: int = 5,
        access_scopes: Optional[List[str]] = None,
        equipment_model: Optional[str] = None,
        strategy: Optional[str] = None,
    ) -> List[Citation]:
        chunks = self._active_chunks(access_scopes)
        selected_strategy = RetrievalStrategy(strategy or self.settings.retrieval_strategy)
        ranked = self.engine.rank(
            query,
            chunks,
            selected_strategy,
            top_k,
            equipment_model=equipment_model,
        )
        citations = []
        for score, chunk, components in ranked:
            metadata = chunk["metadata"]
            citations.append(
                Citation(
                    document_id=chunk["document_id"],
                    chunk_id=chunk["chunk_id"],
                    title=chunk["title"],
                    document_type=chunk["document_type"],
                    chapter=metadata.get("chapter"),
                    page=metadata.get("page"),
                    excerpt=chunk["content"][:700],
                    score=round(score, 4),
                    retrieval_method=selected_strategy.value,
                    score_components=components,
                )
            )
        return citations


def ensure_under_root(candidate: Path, root: Path) -> Path:
    resolved = candidate.resolve()
    resolved_root = root.resolve()
    if resolved != resolved_root and resolved_root not in resolved.parents:
        raise ValueError("source_path 必须位于配置的知识库目录内")
    return resolved
